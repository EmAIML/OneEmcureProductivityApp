import os
import re
import json
import base64
import boto3
from pptx import Presentation
from botocore.exceptions import ClientError
from pptx2txt2 import extract_images
from modules.model2 import main as generate_audio_story

# AWS Bedrock setup
client = boto3.client("bedrock-runtime", region_name="ap-south-1")
MODEL_ID = "anthropic.claude-3-haiku-20240307-v1:0"


def encode_image_to_base64(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def describe_image(image_base64, context_text):
    """Send image and context to Claude for interpretation."""
    prompt = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 1024,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/png",
                            "data": image_base64,
                        },
                    },
                    {
                        "type": "text",
                        "text": context_text,
                    },
                ],
            }
        ],
    }

    body = json.dumps(prompt)
    response = client.invoke_model(modelId=MODEL_ID, body=body)
    data = json.loads(response["body"].read())
    return data["content"][0]["text"]


def clean_text(text):
    """Clean unwanted formatting and metadata from text."""
    text = re.sub(r'https?://\S+', '', text)
    patterns = [
        r'\bstyle\.visibility\b',
        r'\bppt_x\b',
        r'\bppt_y\b',
        r'\bstyle\.visibility\s+\S+',
        r'\bppt_x\s+\S+',
        r'\bppt_y\s+\S+',
    ]
    for pattern in patterns:
        text = re.sub(pattern, '', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()


def extract_text_from_presentation(pptx_path):
    """Extract text content from each slide with proper sequence and flow."""
    prs = Presentation(pptx_path)
    slide_texts = {}
    for idx, slide in enumerate(prs.slides, start=1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        slide_texts[idx] = clean_text(slide_text)
    return slide_texts


def process_pptx(pptx_path, output_txt_path, images_output_dir):
    """Extracts text and image content from slides and describes them."""
    os.makedirs(images_output_dir, exist_ok=True)

    # Extract slide text
    slide_text_map = extract_text_from_presentation(pptx_path)

    # Extract images
    image_paths = extract_images(pptx_path, images_output_dir)

    # Map images to slides
    slide_images_map = {}
    for img_file in sorted(os.listdir(images_output_dir)):
        if img_file.startswith("slide_"):
            parts = img_file.split("_")
            if len(parts) >= 2 and parts[1].isdigit():
                slide_num = int(parts[1])
                slide_images_map.setdefault(slide_num, []).append(img_file)

    print("\n Starting slide analysis...\n")

    # Write extracted data
    with open(output_txt_path, "w", encoding="utf-8") as outf:
        for slide_idx in sorted(slide_text_map.keys()):
            cleaned_text = slide_text_map[slide_idx].strip()
            outf.write(f"\n--- Slide {slide_idx} Text ---\n")
            outf.write((cleaned_text or "[No visible text on this slide]") + "\n")

            # Print slide header
            print(f"\n Slide {slide_idx}:")
            if cleaned_text:
                print(f" Text: {cleaned_text[:120]}{'...' if len(cleaned_text) > 120 else ''}")
            else:
                print(" No text detected.")

            # Process images (always analyze them)
            images_for_slide = slide_images_map.get(slide_idx, [])
            for i, img_file in enumerate(images_for_slide, start=1):
                img_path = os.path.join(images_output_dir, img_file)
                try:
                    image_base64 = encode_image_to_base64(img_path)
                    print(f" Analyzing image {i}/{len(images_for_slide)}: {img_file}")

                    image_prompt = (
    f"Slide {slide_idx} context: \"{cleaned_text}\".\n"
    f"This is image {i} of {len(images_for_slide)} from the slide.\n\n"
    "Analyze this image thoroughly and describe its complete educational meaning.\n"
    "Follow these rules:\n"
    "1. Identify what the image is (table, diagram, chart, flowchart, process, graph, picture, or illustration).\n"
    "2. Extract and describe every visible label, word, heading, number, axis title, arrow direction, or symbol exactly "
    "as shown in the image.\n"
    "3. If the image contains a table, describe the table word-by-word: mention column headers, row headers, and "
    "all cell values clearly and in order.\n"
    "4. If it is a diagram, workflow, cycle, hierarchy, or algorithm, explain each block, step, or node in proper "
    "sequence and describe how they connect.\n"
    "5. If it is a chart or graph, describe what the axes represent, data points shown, trends, and any numerical "
    "information visible.\n"
    "6. If it’s a scientific or technical image (chemical structure, process diagram, organ, machinery, etc.), explain "
    "only what is visually present—do NOT add any external scientific knowledge.\n"
    "7. Do NOT mention slide design, colors, styling, spacing, or layout unless they carry educational meaning.\n"
    "8. Do NOT hallucinate or infer meanings that are not directly visible.\n"
    "Describe ONLY the literal and educational content visible in the image.\n"
)


                    description = describe_image(image_base64, image_prompt)

                except (ClientError, Exception) as e:
                    description = f"ERROR describing image: {e}"

                outf.write(f"\nImage: {img_file}\n")
                outf.write("Description:\n" + description + "\n")

                print(f" Image Description:\n{description}\n")

    print(f"\n Slide analysis complete! Output written to: {output_txt_path}")

    print("\n Generating audio story...\n")
    generate_audio_story(output_txt_path, pause_ms=1500)

    print("\n Story generation complete!\n")
                                                         
if __name__ == "__main__":
    import sys

    if len(sys.argv) != 2:
        print("Usage: python models.py input.pptx")
        sys.exit(1)

    pptx_file = sys.argv[1]

    # Get absolute path of current script
    project_dir = os.path.dirname(os.path.abspath(__file__))

    # Prepare paths for outputs
    pptx_basename = os.path.splitext(os.path.basename(pptx_file))[0]
    output_txt = os.path.join(project_dir, f"{pptx_basename}_output.txt")
    images_dir = os.path.join(project_dir, f"{pptx_basename}_images")

    # Process and generate output
    process_pptx(pptx_file, output_txt, images_dir)
